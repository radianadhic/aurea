"""
AUREA — Gold Standard of Data V1.0
PowerPoint Generator

Based on RKT Global CIF & Product Pricing Architecture v1.00 (Aug 2026)
Adapted to AUREA Master Data Management platform architecture.

Generates: AUREA_Gold_Standard_of_Data_V1.0.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree
from pptx.oxml.ns import qn

def OxmlElement(name):
    return etree.SubElement(etree.Element('root'), qn(name))
from PIL import Image
import copy

# ============================================================
# AUREA BRAND COLORS
# ============================================================
NAVY = RGBColor(0x0A, 0x19, 0x29)
NAVY_LIGHT = RGBColor(0x1A, 0x2F, 0x47)
NAVY_DARK = RGBColor(0x05, 0x0F, 0x19)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
GOLD_LIGHT = RGBColor(0xFF, 0xD7, 0x64)
GOLD_DARK = RGBColor(0xB8, 0x86, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_50 = RGBColor(0xF9, 0xFA, 0xFB)
GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_200 = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_300 = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GRAY_700 = RGBColor(0x37, 0x41, 0x51)
SUCCESS = RGBColor(0x16, 0xA3, 0x4A)
INFO = RGBColor(0x02, 0x84, 0xC7)
WARNING = RGBColor(0xEA, 0x58, 0x0C)

ASSET_DIR = '/home/user/aurea-pptx-assets'

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================
# HELPER FUNCTIONS
# ============================================================
def add_rect(slide, x, y, w, h, fill, line=None, line_w=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_round_rect(slide, x, y, w, h, fill, line=None, line_w=0, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=12, bold=False, color=NAVY, font='Calibri',
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullet_list(slide, x, y, w, h, items, size=14, color=NAVY, font='Calibri',
                    bullet_color=GOLD, line_spacing=1.3):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        # bullet
        r1 = p.add_run()
        r1.text = '◆  '
        r1.font.name = font
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        # text
        r2 = p.add_run()
        r2.text = item
        r2.font.name = font
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_image(slide, filename, x, y, w=None, h=None):
    path = os.path.join(ASSET_DIR, filename)
    if not os.path.exists(path):
        print(f'    ! missing image: {filename}')
        return None
    if w and h:
        return slide.shapes.add_picture(path, x, y, w, h)
    elif w:
        return slide.shapes.add_picture(path, x, y, width=w)
    elif h:
        return slide.shapes.add_picture(path, x, y, height=h)
    else:
        return slide.shapes.add_picture(path, x, y)


def add_image_fit(slide, filename, x, y, max_w, max_h):
    """Add image scaled to fit within max bounds maintaining aspect ratio."""
    path = os.path.join(ASSET_DIR, filename)
    if not os.path.exists(path):
        return None
    img = Image.open(path)
    iw, ih = img.size
    # Max bounds in EMU (Inches)
    max_w_emu = max_w
    max_h_emu = max_h
    # Convert to pixels (96 dpi)
    max_w_px = int(max_w_emu / 914400 * 96)
    max_h_px = int(max_h_emu / 914400 * 96)
    # Scale
    scale = min(max_w_px / iw, max_h_px / ih)
    new_w = int(iw * scale / 96 * 914400)
    new_h = int(ih * scale / 96 * 914400)
    # Center
    cx = x + (max_w - new_w) // 2
    cy = y + (max_h - new_h) // 2
    return slide.shapes.add_picture(path, cx, cy, new_w, new_h)


def add_slide_header(slide, title, subtitle=None, page_num=None, total=16):
    """Add header bar with title for content slide."""
    # Top bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.6), NAVY)
    # Gold accent line
    add_rect(slide, 0, Inches(0.6), SLIDE_W, Inches(0.05), GOLD)
    # AUREA logo on left
    add_text(slide, Inches(0.4), Inches(0.1), Inches(1.5), Inches(0.4),
             'AUREA', size=14, bold=True, color=GOLD, font='Georgia')
    # Page title
    add_text(slide, Inches(2.0), Inches(0.1), Inches(9.0), Inches(0.4),
             title, size=18, bold=True, color=WHITE, font='Georgia')
    if page_num:
        add_text(slide, Inches(11.5), Inches(0.1), Inches(1.7), Inches(0.4),
                 f'{page_num} / {total}', size=10, color=GRAY_300, align=PP_ALIGN.RIGHT)
    # Subtitle below bar
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.4),
                 subtitle, size=12, italic=True, color=GOLD_DARK)
    # Title underline
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(0.6), Inches(0.04), GOLD)


def add_slide_footer(slide):
    """Add footer bar."""
    add_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.25), NAVY)
    add_text(slide, Inches(0.4), Inches(7.28), Inches(8), Inches(0.2),
             'AUREA — The Gold Standard of Data  •  V1.0  •  Bank XYZ Confidential',
             size=8, color=GRAY_300)
    add_text(slide, Inches(8), Inches(7.28), Inches(5.0), Inches(0.2),
             'RKT Global CIF & Product Pricing Architecture',
             size=8, color=GRAY_500, align=PP_ALIGN.RIGHT, italic=True)


# ============================================================
# SLIDE BUILDERS
# ============================================================
def slide_01_cover(prs):
    """Cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Navy background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # Gold accent
    add_rect(slide, 0, Inches(2.0), SLIDE_W, Inches(0.05), GOLD)
    add_rect(slide, 0, Inches(5.5), SLIDE_W, Inches(0.05), GOLD)

    # Diamond mark
    add_text(slide, 0, Inches(0.6), SLIDE_W, Inches(1.0),
             '◆', size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # AUREA
    add_text(slide, 0, Inches(1.7), SLIDE_W, Inches(1.2),
             'AUREA', size=96, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font='Georgia')

    # Tagline
    add_text(slide, 0, Inches(2.9), SLIDE_W, Inches(0.6),
             'THE GOLD STANDARD OF DATA', size=20, bold=True, color=GOLD_LIGHT,
             align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(slide, 0, Inches(3.7), SLIDE_W, Inches(0.5),
             'Global CIF & Product Pricing Architecture', size=18, color=WHITE,
             align=PP_ALIGN.CENTER, italic=True)

    # Description
    add_text(slide, Inches(2), Inches(4.5), Inches(9.333), Inches(0.8),
             'A single customer view and a single commercial answer\n'
             'before any channel creates a contract in core banking.',
             size=15, color=GRAY_300, align=PP_ALIGN.CENTER)

    # Bottom block
    add_text(slide, 0, Inches(5.7), SLIDE_W, Inches(0.5),
             'BANK XYZ  •  ARCHITECTURE CONCEPT  •  AUGUST 2026',
             size=12, bold=True, color=GOLD_DARK, align=PP_ALIGN.CENTER)

    add_text(slide, 0, Inches(6.2), SLIDE_W, Inches(0.4),
             'Adapted to AUREA Master Data Management Platform',
             size=11, color=GRAY_500, align=PP_ALIGN.CENTER, italic=True)

    # Footer
    add_text(slide, 0, Inches(6.8), SLIDE_W, Inches(0.3),
             'PARTY TRUTH  ◆  COMMERCIAL TRUTH  ◆  ONE EXPERIENCE',
             size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_02_agenda(prs):
    """Agenda / TOC."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Agenda', 'The story we will tell today')

    # Two-column agenda
    items = [
        ('01', 'The Problem', 'Why three truths matter'),
        ('02', 'Three Truths, One Experience', 'Design principle'),
        ('03', 'Master Infographic', 'End-to-end view'),
        ('04', 'Customer Journey', '7 steps from Discover to Serve'),
        ('05', 'Global CIF Architecture', '6-step identity resolution'),
        ('06', 'Product Commercialization', 'Catalogue, offer, snapshot'),
        ('07', 'Pricing, Eligibility & Limit', 'One explainable offer'),
        ('08', 'Information Model', 'P → O → A identifier chain'),
    ]
    items2 = [
        ('09', 'Reference Technology Stack', 'Tools that matter'),
        ('10', 'Integration Sequence', 'One booking, all channels'),
        ('11', 'Operating Model', 'Governance as first-class'),
        ('12', 'Implementation Path', '4 waves to value'),
        ('13', 'Data Landing', 'Bronze → Silver → Gold'),
        ('14', 'Enterprise System Coverage', 'Every system participates'),
        ('15', 'Reference Implementation', 'Feasible, open stack'),
        ('16', 'Value & Next Steps', 'Outcomes and call to action'),
    ]

    col_w = Inches(5.8)
    row_h = Inches(0.65)
    start_y = Inches(1.6)

    for col_idx, col_items in enumerate([items, items2]):
        x = Inches(0.6) + col_idx * (col_w + Inches(0.5))
        for i, (num, title, desc) in enumerate(col_items):
            y = start_y + i * row_h
            # Number circle
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.1), Inches(0.45), Inches(0.45))
            circle.fill.solid()
            circle.fill.fore_color.rgb = GOLD
            circle.line.color.rgb = NAVY
            circle.line.width = Pt(1.5)
            circle.shadow.inherit = False
            add_text(slide, x, y + Inches(0.13), Inches(0.45), Inches(0.4),
                     num, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
            # Title
            add_text(slide, x + Inches(0.6), y, col_w - Inches(0.6), Inches(0.3),
                     title, size=12, bold=True, color=NAVY)
            # Description
            add_text(slide, x + Inches(0.6), y + Inches(0.3), col_w - Inches(0.6), Inches(0.3),
                     desc, size=9, color=GRAY_500, italic=True)

    add_slide_footer(slide)


def slide_03_problem(prs):
    """The Problem — why three truths matter."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'The Problem', 'Why "customer" is a moving target across the bank')

    # Left side: 4 problem cards
    problems = [
        {
            'icon': '◈',
            'title': 'NO SINGLE CUSTOMER VIEW',
            'desc': 'Each core, CRM, and channel has its own CIF. The same person is three different records.',
        },
        {
            'icon': '◈',
            'title': 'PRICE & LIMIT RE-CALCULATED AT BOOKING',
            'desc': 'The offer shown to the customer does not match the rate the core books. Disputes and rework.',
        },
        {
            'icon': '◈',
            'title': 'DUPLICATE CUSTOMERS, DUPLICATE LIMITS',
            'desc': 'A "new" prospect is actually an existing client. Exposure is miscalculated across systems.',
        },
        {
            'icon': '◈',
            'title': 'CHANNEL-SPECIFIC EXPERIENCE',
            'desc': 'Mobile, branch, RM, and partner channels all show different offers. No continuity.',
        },
    ]
    cy = Inches(1.55)
    for p in problems:
        # Card
        card = add_round_rect(slide, Inches(0.5), cy, Inches(6.2), Inches(1.15), NAVY, line=GOLD, line_w=1.5, radius=0.05)
        # Icon
        add_text(slide, Inches(0.7), cy + Inches(0.2), Inches(0.6), Inches(0.6),
                 p['icon'], size=28, bold=True, color=GOLD)
        # Title
        add_text(slide, Inches(1.4), cy + Inches(0.15), Inches(5.2), Inches(0.4),
                 p['title'], size=12, bold=True, color=GOLD)
        # Desc
        add_text(slide, Inches(1.4), cy + Inches(0.55), Inches(5.2), Inches(0.55),
                 p['desc'], size=10, color=WHITE)
        cy += Inches(1.27)

    # Right side: insight box
    rb_x = Inches(7.1)
    add_round_rect(slide, rb_x, Inches(1.55), Inches(5.7), Inches(5.0), GOLD, line=GOLD_DARK, line_w=2, radius=0.03)
    add_text(slide, rb_x + Inches(0.3), Inches(1.7), Inches(5.1), Inches(0.5),
             '◆ THE INSIGHT', size=12, bold=True, color=NAVY)
    add_text(slide, rb_x + Inches(0.3), Inches(2.2), Inches(5.1), Inches(0.6),
             'Core banking cannot be replaced overnight.',
             size=18, bold=True, color=NAVY, font='Georgia')
    add_text(slide, rb_x + Inches(0.3), Inches(2.9), Inches(5.1), Inches(2.5),
             'But the customer and the commercial offer must be resolved BEFORE any channel creates a contract.\n\n'
             'AUREA sits before the booking systems and publishes one truth back to every channel — '
             'without forcing legacy cores to share a physical customer number.',
             size=12, color=NAVY)

    # Key principle
    add_round_rect(slide, rb_x + Inches(0.3), Inches(5.5), Inches(5.1), Inches(1.0), NAVY, line=GOLD, line_w=1)
    add_text(slide, rb_x + Inches(0.4), Inches(5.6), Inches(5.0), Inches(0.4),
             'KEY PRINCIPLE', size=9, bold=True, color=GOLD)
    add_text(slide, rb_x + Inches(0.4), Inches(5.9), Inches(5.0), Inches(0.6),
             'Channels consume customer & commercial truth first.\n'
             'Cores remain the booking & ledger authority.',
             size=10, color=WHITE, italic=True)

    add_slide_footer(slide)


def slide_04_three_truths(prs):
    """Three Truths, One Experience — design principle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Three Truths, One Experience',
                     'The design principle that drives every decision in AUREA')

    # Add the three truths image
    add_image_fit(slide, 'three_truths.png', Inches(0.5), Inches(1.5),
                  Inches(12.3), Inches(5.5))

    add_slide_footer(slide)


def slide_05_master_infographic(prs):
    """Master Infographic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Master Infographic',
                     'Every channel receives one customer and one offer')

    add_image_fit(slide, 'master_infographic.png', Inches(0.3), Inches(1.4),
                  Inches(12.7), Inches(5.8))

    add_slide_footer(slide)


def slide_06_customer_journey(prs):
    """Customer Journey — 7 steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Customer Journey',
                     'The customer is recognized before the bank creates a contract')

    add_image_fit(slide, 'customer_journey.png', Inches(0.3), Inches(1.4),
                  Inches(12.7), Inches(5.7))

    add_slide_footer(slide)


def slide_07_cif_architecture(prs):
    """Global CIF Architecture — 6-step pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Global CIF Architecture',
                     'A governed identity-resolution service — one enterprise key, many local CIFs')

    add_image_fit(slide, 'cif_pipeline.png', Inches(0.3), Inches(1.4),
                  Inches(12.7), Inches(5.6))

    add_slide_footer(slide)


def slide_08_product_commercialization(prs):
    """Product Commercialization — 6 layers + 3 attribute groups."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Product Commercialization',
                     'A central catalogue creates channel-ready offers')

    # Left: 6-stage flow
    stages = [
        ('Product family', 'Deposit • Loan • Card • Payment'),
        ('Product specification', 'capabilities • lifecycle • accounting map'),
        ('Market offer', 'brand • market • channel • segment • bundle'),
        ('Commercial rules', 'fees • rates • discounts • eligibility'),
        ('Limit policy', 'pre-limit • exposure • collateral • tenor'),
        ('Offer snapshot', 'immutable terms accepted by customer'),
    ]
    sx = Inches(0.5)
    sy = Inches(1.5)
    sw = Inches(0.4)
    sh = Inches(0.7)
    gap = Inches(0.05)
    for i, (title, desc) in enumerate(stages):
        y = sy + i * (sh + gap)
        # Step number
        add_round_rect(slide, sx, y, sw, sh, GOLD, line=GOLD_DARK, line_w=1)
        add_text(slide, sx, y + Inches(0.1), sw, Inches(0.5),
                 str(i+1), size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # Title
        add_round_rect(slide, sx + sw + Inches(0.1), y, Inches(5.3), sh, NAVY, line=GOLD, line_w=1)
        add_text(slide, sx + sw + Inches(0.2), y + Inches(0.1), Inches(5.0), Inches(0.3),
                 title, size=11, bold=True, color=GOLD)
        add_text(slide, sx + sw + Inches(0.2), y + Inches(0.4), Inches(5.0), Inches(0.3),
                 desc, size=10, color=WHITE)
        # Arrow
        if i < 5:
            ay = y + sh + Inches(0.0)
            ax = sx + sw/2
            draw_arrow(slide, ax, ay, ax, ay + gap, color=GOLD)

    # Right: 3 attribute groups
    rb_x = Inches(6.6)
    add_text(slide, rb_x, Inches(1.4), Inches(6.4), Inches(0.4),
             'CATALOGUE ATTRIBUTES EXPOSED BEFORE ONBOARDING', size=12, bold=True, color=NAVY)

    groups = [
        ('Identity & availability', [
            'Product / offer code', 'Market + jurisdiction',
            'Customer type + segment', 'Channel + campaign',
            'Effective-from / to'
        ], GOLD, NAVY),
        ('Commercial terms', [
            'Currency + balance bands', 'Rate / margin / index',
            'Fees + waivers', 'Tenor + repayment',
            'Bundle + benefit'
        ], NAVY_LIGHT, WHITE),
        ('Risk & fulfillment', [
            'Eligibility criteria', 'Policy / pre-approved limit',
            'Collateral + exposure rules', 'Required documents',
            'Disclosures + SLA'
        ], NAVY, WHITE),
    ]
    gy = Inches(1.9)
    for title, items, bg, fg in groups:
        # Card
        add_round_rect(slide, rb_x, gy, Inches(6.4), Inches(1.4), bg, line=GOLD, line_w=1, radius=0.03)
        add_text(slide, rb_x + Inches(0.2), gy + Inches(0.1), Inches(6.0), Inches(0.3),
                 '◆  ' + title, size=12, bold=True, color=GOLD if fg == WHITE else NAVY)
        # Items in 2 columns
        col_w = Inches(3.0)
        for j, item in enumerate(items):
            col = j % 2
            row = j // 2
            ix = rb_x + Inches(0.2) + col * col_w
            iy = gy + Inches(0.5) + row * Inches(0.27)
            add_text(slide, ix, iy, col_w, Inches(0.25),
                     '• ' + item, size=10, color=fg)
        gy += Inches(1.55)

    # Critical distinction
    add_round_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), GOLD_LIGHT, line=GOLD_DARK, line_w=2)
    add_text(slide, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.5),
             '◆  LIMIT DISTINCTION:  Catalogue/decisioning owns policy and pre-approved limits.  '
             'Lending/core owns the booked facility/account limit and utilization.',
             size=11, bold=True, color=NAVY)

    add_slide_footer(slide)


def draw_arrow(slide, x1, y1, x2, y2, color=GOLD, head_w=8):
    """Draw a line with arrow head."""
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    line.line.color.rgb = color
    line.line.width = Pt(2)
    # Add arrowhead
    ln = line.line._get_or_add_ln()
    tail = OxmlElement('a:tailEnd')
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('h', 'med')
    ln.append(tail)


def slide_09_pricing_limit(prs):
    """Pricing, Eligibility & Limit — decisioning."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Pricing, Eligibility & Limit',
                     'One decision sequence produces one explainable offer')

    # 3 columns: Input / Decision / Output
    col_w = Inches(4.0)
    col_h = Inches(5.4)
    col_y = Inches(1.5)
    gap = Inches(0.15)

    # Input
    add_round_rect(slide, Inches(0.5), col_y, col_w, col_h, NAVY, line=GOLD, line_w=2, radius=0.03)
    add_rect(slide, Inches(0.5), col_y, col_w, Inches(0.5), GOLD)
    add_text(slide, Inches(0.5), col_y + Inches(0.05), col_w, Inches(0.4),
             'INPUT CONTEXT', size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), col_y + Inches(0.6), col_w - Inches(0.2), Inches(0.3),
             'Customer + Request', size=11, bold=True, color=GOLD)
    input_items = [
        'Global Party ID / prospect ID',
        'Customer type + segment',
        'KYC / residency / consent',
        'Relationship + holdings',
        'Total exposure + risk grade',
        'Channel + location + campaign',
        'Amount + tenor + collateral',
        'Requested effective date',
    ]
    iy = col_y + Inches(1.0)
    for it in input_items:
        add_text(slide, Inches(0.7), iy, col_w - Inches(0.4), Inches(0.3),
                 '◆  ' + it, size=10, color=WHITE)
        iy += Inches(0.35)

    # Arrow 1
    draw_arrow(slide, Inches(0.5) + col_w + Inches(0.02), col_y + col_h/2,
               Inches(0.5) + col_w + gap - Inches(0.02), col_y + col_h/2, color=GOLD)

    # Decision
    add_round_rect(slide, Inches(0.5) + col_w + gap, col_y, col_w, col_h, GOLD, line=GOLD_DARK, line_w=2, radius=0.03)
    add_rect(slide, Inches(0.5) + col_w + gap, col_y, col_w, Inches(0.5), NAVY)
    add_text(slide, Inches(0.5) + col_w + gap, col_y + Inches(0.05), col_w, Inches(0.4),
             'VERSIONED DECISION', size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5) + col_w + gap + Inches(0.1), col_y + Inches(0.6), col_w - Inches(0.2), Inches(0.3),
             'Offer Composer', size=11, bold=True, color=NAVY)
    decision_items = [
        'Filter eligible offers',
        'Select base price',
        'Apply segment / campaign adjustments',
        'Calculate policy limit + exposure',
        'Assemble documents + disclosures',
        'Persist inputs, rules, results',
    ]
    iy = col_y + Inches(1.0)
    for i, it in enumerate(decision_items, 1):
        add_text(slide, Inches(0.5) + col_w + gap + Inches(0.2), iy, col_w - Inches(0.4), Inches(0.3),
                 f'{i}.  ' + it, size=10, color=NAVY, bold=(i==1))
        iy += Inches(0.45)

    # DMN note
    add_text(slide, Inches(0.5) + col_w + gap + Inches(0.2), col_y + col_h - Inches(0.7), col_w - Inches(0.4), Inches(0.4),
             'Powered by DMN / RULES / MODELS', size=10, bold=True, color=NAVY, italic=True, align=PP_ALIGN.CENTER)

    # Arrow 2
    draw_arrow(slide, Inches(0.5) + 2*col_w + gap + Inches(0.02), col_y + col_h/2,
               Inches(0.5) + 2*col_w + 2*gap - Inches(0.02), col_y + col_h/2, color=GOLD)

    # Output
    ox = Inches(0.5) + 2*col_w + 2*gap
    add_round_rect(slide, ox, col_y, col_w, col_h, NAVY, line=GOLD, line_w=2, radius=0.03)
    add_rect(slide, ox, col_y, col_w, Inches(0.5), GOLD)
    add_text(slide, ox, col_y + Inches(0.05), col_w, Inches(0.4),
             'CONSISTENT OUTPUT', size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(slide, ox + Inches(0.1), col_y + Inches(0.6), col_w - Inches(0.2), Inches(0.3),
             'Channel-ready Answer', size=11, bold=True, color=GOLD)
    output_items = [
        'Eligible products / bundles',
        'Personalized rate and fees',
        'Approved or indicative limit',
        'Expiry and conditions',
        'Required evidence',
        'Disclosures + consent text',
        'Decline / referral reason codes',
        'Offer Snapshot ID',
    ]
    iy = col_y + Inches(1.0)
    for it in output_items:
        add_text(slide, ox + Inches(0.2), iy, col_w - Inches(0.4), Inches(0.3),
                 '◆  ' + it, size=10, color=WHITE)
        iy += Inches(0.35)

    # Bottom callout
    add_round_rect(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.0), NAVY)
    add_text(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.25),
             '◆  No CIF yet?  Create a prospect Party ID, then merge safely after identity verification.',
             size=10, bold=True, color=GOLD_DARK, align=PP_ALIGN.CENTER, italic=True)

    add_slide_footer(slide)


def slide_10_information_model(prs):
    """Information Model — 3 identifiers chain."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Information Model',
                     'Three identifiers keep channels and cores loosely coupled')

    add_image_fit(slide, 'identifier_chain.png', Inches(0.5), Inches(1.4),
                  Inches(12.3), Inches(5.5))

    add_slide_footer(slide)


def slide_11_tech_stack(prs):
    """Reference Technology Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Reference Technology Stack',
                     'Capability boundaries and contracts matter more than any single vendor')

    # 6 layer rows
    layers = [
        ('EXPERIENCE', NAVY_LIGHT, [
            ('Mobile / Web / Branch / CRM', 'Channel UIs and assisted journeys'),
            ('BFF + Design System', 'Channel capability contracts, journey state, accessibility, telemetry'),
        ]),
        ('ACCESS + SECURITY', NAVY, [
            ('API Gateway: Kong / Apigee / WSO2', 'FAPI 2.0 • OIDC • mTLS • Consent'),
            ('Service Mesh + Vault/KMS', 'Zero-trust segmentation, secret management'),
        ]),
        ('JOURNEY + DECISION', GOLD_DARK, [
            ('BPMN: Camunda / Temporal', 'Orchestration, case management'),
            ('DMN / Drools / Offer Composer', 'Decisioning, pricing rules, eligibility'),
        ]),
        ('DOMAIN SERVICES', NAVY, [
            ('Global CIF / MDM + Customer 360', 'Identity, KYC, relationships, segment'),
            ('Product Catalogue + Pricing + Limit', 'Specifications, offers, eligibility, limit'),
        ]),
        ('DATA + EVENTS', NAVY_LIGHT, [
            ('Kafka / Pulsar + Schema Registry', 'Event backbone, schema governance'),
            ('Debezium CDC + PostgreSQL + Redis + OpenSearch', 'Operational data, search, cache'),
        ]),
        ('PLATFORM + CONTROL', NAVY, [
            ('Kubernetes / OpenShift', 'Container orchestration, hybrid / multi-cloud'),
            ('OpenTelemetry + CI/CD + Policy-as-Code', 'Observability, deployment, compliance'),
        ]),
    ]

    ly = Inches(1.5)
    lh = Inches(0.85)
    gap = Inches(0.05)
    for i, (name, color, items) in enumerate(layers):
        # Layer bar
        add_rect(slide, Inches(0.5), ly, Inches(2.5), lh, color)
        add_text(slide, Inches(0.5), ly + Inches(0.25), Inches(2.5), Inches(0.4),
                 name, size=11, bold=True, color=GOLD if color == NAVY else WHITE,
                 align=PP_ALIGN.CENTER)

        # Items
        for j, (tech, desc) in enumerate(items):
            ix = Inches(3.1) + j * Inches(4.7)
            add_round_rect(slide, ix, ly, Inches(4.5), lh, WHITE, line=color, line_w=1, radius=0.05)
            add_text(slide, ix + Inches(0.15), ly + Inches(0.1), Inches(4.2), Inches(0.3),
                     tech, size=10, bold=True, color=color)
            add_text(slide, ix + Inches(0.15), ly + Inches(0.4), Inches(4.2), Inches(0.4),
                     desc, size=9, color=GRAY_700)
        ly += lh + gap

    # Bottom: deployment notes
    add_round_rect(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.25), NAVY, line=GOLD, line_w=1)
    add_text(slide, Inches(0.7), Inches(7.03), Inches(12.0), Inches(0.2),
             '◆  Hybrid / on-prem / cloud   •   Active-active where justified   •   Zero-trust segmentation   •   Immutable audit   •   Data-residency controls',
             size=9, bold=True, color=GOLD_LIGHT, italic=True)

    add_slide_footer(slide)


def slide_12_integration_sequence(prs):
    """Integration Sequence — 8 steps with swimlanes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Integration Sequence',
                     'Core books the same answer every channel presented')

    add_image_fit(slide, 'integration_sequence.png', Inches(0.3), Inches(1.4),
                  Inches(12.7), Inches(5.6))

    add_slide_footer(slide)


def slide_13_operating_model(prs):
    """Operating Model — governance, lifecycle, roles."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Operating Model',
                     'Governance keeps customer and commercial truth reliable')

    add_image_fit(slide, 'operating_model.png', Inches(0.3), Inches(1.4),
                  Inches(12.7), Inches(5.6))

    add_slide_footer(slide)


def slide_14_implementation_path(prs):
    """Implementation Path — 4 waves."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Implementation Path',
                     'Deliver value in waves without a core replacement')

    add_image_fit(slide, 'implementation_path.png', Inches(0.3), Inches(1.4),
                  Inches(12.7), Inches(5.7))

    add_slide_footer(slide)


def slide_15_data_landing(prs):
    """Data Landing — Bronze → Silver → Gold."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Data Landing',
                     'Operational truth resolves identity; the lakehouse standardizes events')

    add_image_fit(slide, 'data_landing.png', Inches(0.3), Inches(1.4),
                  Inches(12.7), Inches(5.7))

    add_slide_footer(slide)


def slide_16_system_coverage(prs):
    """Enterprise System Coverage."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Enterprise System Coverage',
                     'Every banking system participates through explicit contracts')

    # 3 columns
    # Initiators
    add_text(slide, Inches(0.5), Inches(1.4), Inches(3.8), Inches(0.4),
             'INITIATORS + EXTERNAL DATA', size=11, bold=True, color=NAVY)
    initiators = [
        ('M', 'Mobile'),
        ('W', 'Web'),
        ('BR', 'Branch'),
        ('RM', 'Relationship Manager'),
        ('CC', 'Contact Centre'),
        ('API', 'Partner / Open API'),
    ]
    iy = Inches(1.9)
    for code, name in initiators:
        # Code pill
        add_round_rect(slide, Inches(0.5), iy, Inches(0.7), Inches(0.4), NAVY, line=GOLD, line_w=1, radius=0.2)
        add_text(slide, Inches(0.5), iy + Inches(0.05), Inches(0.7), Inches(0.3),
                 code, size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        # Name
        add_text(slide, Inches(1.3), iy + Inches(0.05), Inches(3.0), Inches(0.3),
                 name, size=11, color=NAVY)
        iy += Inches(0.45)

    add_text(slide, Inches(0.5), Inches(4.7), Inches(3.8), Inches(0.4),
             'MARKET / REFERENCE FEEDS', size=10, bold=True, color=NAVY)
    feeds = ['Bloomberg B-PIPE + Data License', 'LSEG Real-Time (Refinitiv)', 'Exchanges / FX / curves', 'Ratings / sanctions / bureau']
    iy = Inches(5.1)
    for f in feeds:
        add_text(slide, Inches(0.5), iy, Inches(3.8), Inches(0.3),
                 '◆  ' + f, size=9, color=GRAY_700)
        iy += Inches(0.3)

    # Middle: Contract Layer
    cx = Inches(4.7)
    add_text(slide, cx, Inches(1.4), Inches(3.8), Inches(0.4),
             'ENTERPRISE CONTRACT LAYER', size=11, bold=True, color=NAVY)

    add_round_rect(slide, cx, Inches(1.9), Inches(3.8), Inches(2.5), NAVY, line=GOLD, line_w=2, radius=0.03)
    contracts = [
        ('Customer + Offer Façade', 'resolveCustomer • getContext\nlistEligibleOffers • calculatePriceAndLimit\nstartOnboarding • createBookingRequest'),
        ('Global CIF', 'GID + crosswalk • KYC + consent\nrelationships'),
        ('Product / Price / Limit', 'catalogue + rules\nmarket adjustment • offer snapshot'),
    ]
    cy = Inches(2.05)
    for name, desc in contracts:
        add_text(slide, cx + Inches(0.2), cy, Inches(3.5), Inches(0.3),
                 '◆  ' + name, size=10, bold=True, color=GOLD)
        add_text(slide, cx + Inches(0.4), cy + Inches(0.3), Inches(3.4), Inches(0.4),
                 desc, size=8, color=WHITE)
        cy += Inches(0.8)

    # Integration patterns
    add_round_rect(slide, cx, Inches(4.55), Inches(3.8), Inches(1.0), GOLD, line=GOLD_DARK, line_w=1, radius=0.05)
    add_text(slide, cx + Inches(0.1), Inches(4.6), Inches(3.6), Inches(0.3),
             'INTEGRATION PATTERNS', size=9, bold=True, color=NAVY)
    add_text(slide, cx + Inches(0.1), Inches(4.9), Inches(3.6), Inches(0.6),
             '◆  API / synchronous  ◆  Events / async\n'
             '◆  CDC / database log  ◆  Batch / file\n'
             'Schema registry • idempotency • lineage',
             size=8, color=NAVY)

    # Right: Banking sources
    rx = Inches(8.9)
    add_text(slide, rx, Inches(1.4), Inches(3.9), Inches(0.4),
             'BANKING SOURCES + BOOKING SYSTEMS', size=11, bold=True, color=NAVY)

    systems = [
        ('CB', 'Core Banking', 'customer / deposit / ledger'),
        ('TR', 'Treasury', 'deals / positions / liquidity'),
        ('TF', 'Trade Finance', 'LC / guarantee / party'),
        ('CM', 'Cash Management', 'cash products / mandates'),
        ('CD', 'Card System', 'cardholder / account / auth'),
        ('LN', 'Lending', 'facility / collateral / exposure'),
        ('PY', 'Payments', 'instruction / beneficiary / status'),
        ('CX', 'CRM / KYC', 'lead / interaction / CDD'),
        ('FA', 'Fraud / AML', 'alert / score / disposition'),
        ('GL', 'GL / Finance', 'posting / balance / reporting'),
        ('DW', 'Digital Banking', 'session / journey / consent'),
        ('MD', 'Market Data Hub', 'quotes / curves / entitlement'),
    ]
    iy = Inches(1.9)
    col_w = Inches(1.9)
    for i, (code, name, desc) in enumerate(systems):
        col = i % 2
        row = i // 2
        x = rx + col * col_w
        y = iy + row * Inches(0.45)
        # Code pill
        add_round_rect(slide, x, y, Inches(0.4), Inches(0.35), GOLD, line=GOLD_DARK, line_w=1, radius=0.2)
        add_text(slide, x, y + Inches(0.04), Inches(0.4), Inches(0.3),
                 code, size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        # Name
        add_text(slide, x + Inches(0.45), y, col_w - Inches(0.45), Inches(0.2),
                 name, size=8, bold=True, color=NAVY)
        # Desc
        add_text(slide, x + Inches(0.45), y + Inches(0.18), col_w - Inches(0.45), Inches(0.2),
                 desc, size=7, color=GRAY_500)

    # Bottom flow
    add_round_rect(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.55), NAVY, line=GOLD, line_w=1)
    add_text(slide, Inches(0.7), Inches(6.78), Inches(12.0), Inches(0.4),
             'ALL SYSTEMS  →  Kafka/events, Debezium/CDC, controlled batch/files  →  Raw landing  →  '
             'Conformance  →  Global CIF / product / account facts  →  DWH + data products',
             size=10, bold=True, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

    add_slide_footer(slide)


def slide_17_reference_implementation(prs):
    """Reference Implementation Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Reference Implementation Stack',
                     'A feasible stack supports fast decisions and open analytical data')

    add_image_fit(slide, 'reference_stack.png', Inches(0.3), Inches(1.4),
                  Inches(12.7), Inches(5.7))

    add_slide_footer(slide)


def slide_18_value_benefits(prs):
    """Value & Benefits — outcomes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Business Value & Outcomes',
                     'What the bank gains from one truth and one offer everywhere')

    add_image_fit(slide, 'value_benefits.png', Inches(0.3), Inches(1.4),
                  Inches(12.7), Inches(5.7))

    add_slide_footer(slide)


def slide_19_next_steps(prs):
    """Next Steps — call to action."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, 'Next Steps',
                     'How we turn the concept into a working platform')

    # 3 columns
    items = [
        {
            'num': '01',
            'title': 'DECIDE',
            'subtitle': 'Now',
            'color': NAVY,
            'items': [
                'Executive sponsor confirmed',
                'AUREA platform scope agreed',
                'Wave 0 funding approved',
                'Source-system assessment kickoff',
            ]
        },
        {
            'num': '02',
            'title': 'BUILD',
            'subtitle': 'Next 16 weeks',
            'color': GOLD,
            'items': [
                'Wave 0: foundation + governance',
                'Wave 1: Global Party ID + Customer Context API',
                'Stewardship, DQ thresholds, consent',
                'Sandbox channel for early validation',
            ]
        },
        {
            'num': '03',
            'title': 'SCALE',
            'subtitle': 'Beyond',
            'color': NAVY_LIGHT,
            'items': [
                'Wave 2: catalogue, pricing, Offer Snapshot API',
                'Wave 3: priority journeys + core adapters',
                'Event-driven analytics, 360 productization',
                'Open Banking partner API',
            ]
        },
    ]

    col_w = Inches(4.0)
    col_h = Inches(4.0)
    col_y = Inches(1.5)
    gap = Inches(0.15)
    for i, it in enumerate(items):
        x = Inches(0.5) + i * (col_w + gap)
        add_round_rect(slide, x, col_y, col_w, col_h, it['color'], line=GOLD, line_w=2, radius=0.03)
        # Number
        add_text(slide, x + Inches(0.3), col_y + Inches(0.2), col_w - Inches(0.6), Inches(0.6),
                 it['num'], size=42, bold=True, color=GOLD, font='Georgia')
        # Title
        text_fill = GOLD if it['color'] != GOLD else NAVY
        add_text(slide, x + Inches(0.3), col_y + Inches(0.9), col_w - Inches(0.6), Inches(0.4),
                 it['title'], size=22, bold=True, color=text_fill, font='Georgia')
        # Subtitle
        sub_fill = GRAY_300 if it['color'] == NAVY else GRAY_700
        if it['color'] == GOLD:
            sub_fill = NAVY
        add_text(slide, x + Inches(0.3), col_y + Inches(1.4), col_w - Inches(0.6), Inches(0.3),
                 it['subtitle'], size=12, bold=True, color=sub_fill)
        # Divider
        add_rect(slide, x + Inches(0.3), col_y + Inches(1.8), Inches(0.6), Inches(0.03), GOLD)
        # Items
        iy = col_y + Inches(2.0)
        for it_text in it['items']:
            add_text(slide, x + Inches(0.3), iy, col_w - Inches(0.6), Inches(0.4),
                     '◆  ' + it_text, size=11, color=WHITE if it['color'] != GOLD else NAVY)
            iy += Inches(0.45)

    # Bottom: closing
    add_round_rect(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.4), NAVY, line=GOLD, line_w=2)
    add_text(slide, Inches(0.7), Inches(5.85), Inches(12.0), Inches(0.5),
             '◆  THE INVITATION', size=12, bold=True, color=GOLD)
    add_text(slide, Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.9),
             'AUREA is not a new core. It is the layer that lets every channel and every core\n'
             'see the same customer and the same offer — at last.',
             size=18, color=WHITE, italic=True, font='Georgia', align=PP_ALIGN.CENTER)

    add_slide_footer(slide)


def slide_20_thank_you(prs):
    """Thank you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    # Gold accents
    add_rect(slide, 0, Inches(2.5), SLIDE_W, Inches(0.05), GOLD)
    add_rect(slide, 0, Inches(5.0), SLIDE_W, Inches(0.05), GOLD)

    # Diamond
    add_text(slide, 0, Inches(0.6), SLIDE_W, Inches(1.0),
             '◆', size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Thank you
    add_text(slide, 0, Inches(1.8), SLIDE_W, Inches(1.0),
             'THANK YOU', size=80, bold=True, color=GOLD, font='Georgia', align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(slide, 0, Inches(3.0), SLIDE_W, Inches(0.5),
             'AUREA — The Gold Standard of Data', size=20, color=GOLD_LIGHT,
             align=PP_ALIGN.CENTER, italic=True)

    add_text(slide, 0, Inches(3.5), SLIDE_W, Inches(0.5),
             'Global CIF & Product Pricing Architecture', size=14, color=WHITE,
             align=PP_ALIGN.CENTER)

    # Three truths reminder
    add_text(slide, 0, Inches(4.2), SLIDE_W, Inches(0.5),
             'PARTY TRUTH  ◆  COMMERCIAL TRUTH  ◆  ONE EXPERIENCE',
             size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Contact
    add_text(slide, 0, Inches(5.3), SLIDE_W, Inches(0.4),
             'BANK XYZ  •  DATA PLATFORM ENGINEERING  •  AUGUST 2026',
             size=12, bold=True, color=GOLD_DARK, align=PP_ALIGN.CENTER)

    add_text(slide, 0, Inches(5.7), SLIDE_W, Inches(0.4),
             'Adapted from RKT Global CIF & Product Pricing Architecture v1.00',
             size=10, color=GRAY_500, align=PP_ALIGN.CENTER, italic=True)

    add_text(slide, 0, Inches(6.5), SLIDE_W, Inches(0.4),
             'CONFIDENTIAL  •  FOR INTERNAL DISCUSSION',
             size=9, color=GRAY_500, align=PP_ALIGN.CENTER, italic=True)


# ============================================================
# MAIN
# ============================================================
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print('Building AUREA Gold Standard of Data V1.0 presentation...')

    # Cover
    slide_01_cover(prs); print('  ✓ Slide 1: Cover')
    # Agenda
    slide_02_agenda(prs); print('  ✓ Slide 2: Agenda')
    # Problem
    slide_03_problem(prs); print('  ✓ Slide 3: The Problem')
    # Three Truths
    slide_04_three_truths(prs); print('  ✓ Slide 4: Three Truths, One Experience')
    # Master Infographic
    slide_05_master_infographic(prs); print('  ✓ Slide 5: Master Infographic')
    # Customer Journey
    slide_06_customer_journey(prs); print('  ✓ Slide 6: Customer Journey')
    # CIF Architecture
    slide_07_cif_architecture(prs); print('  ✓ Slide 7: Global CIF Architecture')
    # Product Commercialization
    slide_08_product_commercialization(prs); print('  ✓ Slide 8: Product Commercialization')
    # Pricing & Limit
    slide_09_pricing_limit(prs); print('  ✓ Slide 9: Pricing, Eligibility & Limit')
    # Information Model
    slide_10_information_model(prs); print('  ✓ Slide 10: Information Model')
    # Tech Stack
    slide_11_tech_stack(prs); print('  ✓ Slide 11: Reference Tech Stack')
    # Integration Sequence
    slide_12_integration_sequence(prs); print('  ✓ Slide 12: Integration Sequence')
    # Operating Model
    slide_13_operating_model(prs); print('  ✓ Slide 13: Operating Model')
    # Implementation Path
    slide_14_implementation_path(prs); print('  ✓ Slide 14: Implementation Path')
    # Data Landing
    slide_15_data_landing(prs); print('  ✓ Slide 15: Data Landing')
    # System Coverage
    slide_16_system_coverage(prs); print('  ✓ Slide 16: Enterprise System Coverage')
    # Reference Implementation
    slide_17_reference_implementation(prs); print('  ✓ Slide 17: Reference Implementation')
    # Value & Benefits
    slide_18_value_benefits(prs); print('  ✓ Slide 18: Value & Benefits')
    # Next Steps
    slide_19_next_steps(prs); print('  ✓ Slide 19: Next Steps')
    # Thank You
    slide_20_thank_you(prs); print('  ✓ Slide 20: Thank You')

    output = '/home/user/AUREA_Gold_Standard_of_Data_V1.0.pptx'
    prs.save(output)
    size_kb = os.path.getsize(output) / 1024
    print(f'\n✓ PPTX created: {output}')
    print(f'  Size: {size_kb:.0f} KB ({size_kb/1024:.2f} MB)')
    print(f'  Slides: 20')


if __name__ == '__main__':
    main()
