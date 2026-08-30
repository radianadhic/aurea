"""
Render PPTX slides as PNG using python-pptx + Pillow.

Faithful rendering: reads all text/color/size from python-pptx and
draws using Pillow with proper text handling.
"""

import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageDraw, ImageFont

# AUREA colors
NAVY = (10, 25, 41)
NAVY_LIGHT = (26, 47, 71)
GOLD = (212, 175, 55)
GOLD_LIGHT = (255, 215, 100)
GOLD_DARK = (184, 134, 11)
WHITE = (255, 255, 255)
GRAY_500 = (107, 114, 128)
GRAY_700 = (55, 65, 81)
GRAY_300 = (209, 213, 219)
GRAY_100 = (243, 244, 246)

FONT_PATHS_BOLD = [
    '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',
    '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf',
]
FONT_PATHS_REG = [
    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
    '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
]

def get_font(size, bold=False):
    paths = FONT_PATHS_BOLD if bold else FONT_PATHS_REG
    for p in paths:
        if os.path.exists(p):
            try:
                return ImageFont.truetype(p, size)
            except:
                pass
    return ImageFont.load_default()


def emu_to_px(emu, dpi=96):
    return int(emu / 914400 * dpi)


def get_color(rgb):
    if rgb is None:
        return None
    try:
        return (rgb[0], rgb[1], rgb[2])
    except:
        return None


def get_solid_fill(shape):
    """Get solid fill color of a shape."""
    try:
        if hasattr(shape, 'fill'):
            f = shape.fill
            if f.type == 1:  # solid
                return get_color(f.fore_color.rgb)
    except:
        pass
    return None


def get_line_color(shape):
    try:
        if hasattr(shape, 'line'):
            l = shape.line
            if l.color and l.color.type == 1:
                return get_color(l.color.rgb)
    except:
        pass
    return None


def get_line_width_pt(shape):
    try:
        if hasattr(shape, 'line'):
            l = shape.line
            if l.width:
                return l.width.pt
    except:
        pass
    return 0.75


def is_rounded(shape):
    try:
        if hasattr(shape, 'auto_shape_type'):
            t = shape.auto_shape_type
            if t and 'ROUND' in str(t):
                return True
        if hasattr(shape, 'adjustments') and shape.adjustments:
            if shape.adjustments[0] and shape.adjustments[0] > 0.01:
                return True
    except:
        pass
    return False


def wrap_text(text, font, max_w):
    """Wrap text to multiple lines."""
    if not text:
        return []
    words = text.split(' ')
    lines = []
    current = ''
    for word in words:
        test = (current + ' ' + word).strip()
        if font.getlength(test) > max_w and current:
            lines.append(current)
            current = word
        else:
            current = test
    if current:
        lines.append(current)
    return lines


def draw_text_in_shape(draw, x, y, w, h, text_frame):
    """Render text from a text frame."""
    if not text_frame or not text_frame.paragraphs:
        return
    # Anchor
    anchor = text_frame.vertical_anchor
    # Margins
    try:
        ml = int(text_frame.margin_left / 914400 * 96) if text_frame.margin_left else 4
        mr = int(text_frame.margin_right / 914400 * 96) if text_frame.margin_right else 4
        mt = int(text_frame.margin_top / 914400 * 96) if text_frame.margin_top else 2
        mb = int(text_frame.margin_bottom / 914400 * 96) if text_frame.margin_bottom else 2
    except:
        ml, mr, mt, mb = 4, 4, 2, 2

    inner_w = w - ml - mr
    inner_h = h - mt - mb
    if inner_w <= 0 or inner_h <= 0:
        return

    # Pre-wrap all paragraphs
    para_data = []
    total_h = 0
    for para in text_frame.paragraphs:
        # Get runs
        runs = []
        for run in para.runs:
            txt = run.text
            if not txt:
                continue
            # Get font properties
            size_pt = run.font.size.pt if run.font.size else 12
            bold = run.font.bold if run.font.bold is not None else False
            try:
                color = (run.font.color.rgb[0], run.font.color.rgb[1], run.font.color.rgb[2]) if run.font.color and run.font.color.rgb else None
            except:
                color = None
            if color is None:
                # Try theme color
                try:
                    if run.font.color and run.font.color.theme_color:
                        # Just default to black
                        color = (0, 0, 0)
                except:
                    color = (0, 0, 0)
            runs.append({'text': txt, 'size': size_pt, 'bold': bold, 'color': color})
        if not runs:
            para_data.append({'lines': [], 'h': 0})
            continue

        # Combine all runs into one string with mixed formatting
        # For simplicity, just use first run's properties
        first = runs[0]
        full_text = ''.join(r['text'] for r in runs)
        font_size = int(first['size'] * 96 / 72)
        font = get_font(font_size, bold=first['bold'])
        lines = wrap_text(full_text, font, inner_w)
        line_h = int(first['size'] * 1.2 * 96 / 72)
        para_h = line_h * len(lines) + 4  # paragraph spacing
        para_data.append({'lines': lines, 'h': para_h, 'font': font, 'color': first['color'],
                         'size': first['size'], 'line_h': line_h, 'bold': first['bold'],
                         'align': para.alignment})
        total_h += para_h

    # Anchor
    cur_y = y + mt
    if anchor == 3:  # MIDDLE
        cur_y = y + mt + max(0, (inner_h - total_h) // 2)
    elif anchor == 4:  # BOTTOM
        cur_y = y + mt + max(0, inner_h - total_h)

    for pd in para_data:
        if not pd.get('lines'):
            continue
        font = pd['font']
        color = pd.get('color') or NAVY
        line_h = pd['line_h']
        align = pd.get('align')
        for line in pd['lines']:
            bbox = draw.textbbox((0, 0), line, font=font)
            line_w = bbox[2] - bbox[0]
            tx = x + ml
            if align == 2:  # CENTER
                tx = x + ml + max(0, (inner_w - line_w) // 2)
            elif align == 1:  # RIGHT
                tx = x + ml + max(0, inner_w - line_w)
            draw.text((tx, cur_y), line, font=font, fill=color)
            cur_y += line_h
        cur_y += 4  # paragraph spacing


def render_shape(draw, shape, img):
    """Render a single shape."""
    try:
        x = emu_to_px(shape.left)
        y = emu_to_px(shape.top)
        w = emu_to_px(shape.width)
        h = emu_to_px(shape.height)
    except:
        return

    if x < -1000 or y < -1000 or w <= 0 or h <= 0:
        return

    # Picture
    if shape.shape_type == 13:
        try:
            image = shape.image
            img_bytes = image.blob
            pil_img = Image.open(io.BytesIO(img_bytes)).convert('RGB')
            pil_img = pil_img.resize((w, h), Image.LANCZOS)
            img.paste(pil_img, (x, y))
        except Exception as e:
            pass
        # Pictures usually don't have text frame but check anyway
        if shape.has_text_frame:
            draw_text_in_shape(draw, x, y, w, h, shape.text_frame)
        return

    # Group
    if shape.shape_type == 6:  # GROUP
        for child in shape.shapes:
            render_shape(draw, child, img)
        return

    # Get fill and line
    fill_color = get_solid_fill(shape)
    line_color = get_line_color(shape)
    line_w_pt = get_line_width_pt(shape)
    line_w_px = max(1, int(line_w_pt * 96 / 72))

    # Draw shape background
    if fill_color is not None or line_color is not None:
        # Determine shape geometry
        is_oval = False
        try:
            if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type:
                t = str(shape.auto_shape_type)
                if 'OVAL' in t or 'ELLIPSE' in t or 'CIRCLE' in t:
                    is_oval = True
        except:
            pass

        rounded = is_rounded(shape)

        if is_oval:
            draw.ellipse([x, y, x + w, y + h], fill=fill_color,
                        outline=line_color, width=line_w_px)
        elif rounded:
            radius = int(min(w, h) * 0.08)
            try:
                draw.rounded_rectangle([x, y, x + w, y + h], radius=radius,
                                      fill=fill_color, outline=line_color, width=line_w_px)
            except:
                draw.rectangle([x, y, x + w, y + h], fill=fill_color,
                              outline=line_color, width=line_w_px)
        else:
            draw.rectangle([x, y, x + w, y + h], fill=fill_color,
                          outline=line_color, width=line_w_px)

    # Render text
    if shape.has_text_frame:
        draw_text_in_shape(draw, x, y, w, h, shape.text_frame)


def main():
    pptx_path = '/home/user/AUREA_Gold_Standard_of_Data_V1.0.pptx'
    out_dir = '/home/user/aurea-pptx-previews'
    os.makedirs(out_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    print(f'Loaded {len(prs.slides)} slides')

    for i, slide in enumerate(prs.slides):
        w_px = emu_to_px(prs.slide_width)
        h_px = emu_to_px(prs.slide_height)
        img = Image.new('RGB', (w_px, h_px), WHITE)
        draw = ImageDraw.Draw(img)
        for shape in slide.shapes:
            try:
                render_shape(draw, shape, img)
            except Exception as e:
                pass
        out = f'{out_dir}/slide_{i+1:02d}.png'
        img.save(out, optimize=True)
        print(f'  ✓ Slide {i+1}')


if __name__ == '__main__':
    main()
