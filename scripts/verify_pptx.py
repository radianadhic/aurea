"""
Render PPTX slides as PNG by using python-pptx to read the file
and create a faithful rendering using Pillow.

This is a custom render that mimics how PowerPoint will display the slides.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import io

# AUREA colors
NAVY = (10, 25, 41)
NAVY_LIGHT = (26, 47, 71)
NAVY_DARK = (5, 15, 25)
GOLD = (212, 175, 55)
GOLD_LIGHT = (255, 215, 100)
GOLD_DARK = (184, 134, 11)
WHITE = (255, 255, 255)
GRAY_50 = (249, 250, 251)
GRAY_100 = (243, 244, 246)
GRAY_200 = (229, 231, 235)
GRAY_300 = (209, 213, 219)
GRAY_500 = (107, 114, 128)
GRAY_700 = (55, 65, 81)
SUCCESS = (22, 163, 74)
INFO = (2, 132, 199)
WARNING = (234, 88, 12)

FONT_PATHS = [
    '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',
    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
    '/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf',
    '/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf',
]

def get_font(size, bold=False, italic=False):
    paths = []
    if bold:
        paths.extend([p for p in FONT_PATHS if 'Bold' in p])
    else:
        paths.extend([p for p in FONT_PATHS if 'Bold' not in p])
    paths.extend(FONT_PATHS)
    for p in paths:
        if os.path.exists(p):
            try:
                f = ImageFont.truetype(p, size)
                return f
            except:
                pass
    return ImageFont.load_default()


def emu_to_px(emu, dpi=96):
    return int(emu / 914400 * dpi)


def get_color(rgb_color):
    """Convert RGBColor to tuple."""
    if rgb_color is None:
        return None
    try:
        return (rgb_color[0], rgb_color[1], rgb_color[2])
    except:
        return NAVY


def render_text(draw, x, y, w, h, text_frame, default_size=12):
    """Render a text frame."""
    if not text_frame.paragraphs:
        return
    # Vertical anchor
    anchor = text_frame.vertical_anchor
    # Calculate total height
    paragraphs = text_frame.paragraphs
    para_heights = []
    for para in paragraphs:
        # Estimate height
        runs = para.runs
        if not runs:
            para_heights.append(20)
            continue
        max_size = 0
        for run in runs:
            size_pt = run.font.size.pt if run.font.size else 12
            max_size = max(max_size, size_pt)
        line_h = int(max_size * 1.3 * 96 / 72)
        n_lines = max(1, sum(len(r.text) // 60 + 1 for r in runs))
        para_heights.append(line_h * n_lines + 4)
    total_h = sum(para_heights)
    cur_y = y
    if anchor == 1:  # middle
        cur_y = y + (h - total_h) // 2
    elif anchor == 2:  # bottom
        cur_y = y + h - total_h

    for para, p_h in zip(paragraphs, para_heights):
        align = para.alignment
        # Get full text
        text = ''.join(r.text for r in para.runs)
        if not text:
            cur_y += p_h
            continue
        # Get first run for styling
        run0 = para.runs[0] if para.runs else None
        size_pt = run0.font.size.pt if run0 and run0.font.size else default_size
        bold = run0.font.bold if run0 else False
        italic = run0.font.italic if run0 else False
        try:
            color = (run0.font.color.rgb[0], run0.font.color.rgb[1], run0.font.color.rgb[2]) if run0 and run0.font.color and run0.font.color.type else NAVY
        except:
            color = NAVY

        font = get_font(int(size_pt * 96 / 72), bold=bold, italic=italic)

        # Word wrap
        words = text.split(' ')
        lines = []
        current = ''
        for word in words:
            test = (current + ' ' + word).strip()
            if font.getlength(test) > w - 4:
                if current:
                    lines.append(current)
                current = word
            else:
                current = test
        if current:
            lines.append(current)

        line_h = int(size_pt * 1.3 * 96 / 72)
        for line in lines:
            bbox = draw.textbbox((0, 0), line, font=font)
            line_w = bbox[2] - bbox[0]
            if align == 2:  # center
                draw.text((x + (w - line_w) // 2, cur_y), line, font=font, fill=color)
            elif align == 1:  # right
                draw.text((x + w - line_w - 2, cur_y), line, font=font, fill=color)
            else:  # left
                draw.text((x + 2, cur_y), line, font=font, fill=color)
            cur_y += line_h
        cur_y += 4


def render_shape(draw, shape, img):
    """Render a single shape."""
    try:
        x = emu_to_px(shape.left)
        y = emu_to_px(shape.top)
        w = emu_to_px(shape.width)
        h = emu_to_px(shape.height)
    except:
        return

    if x < 0 or y < 0 or w <= 0 or h <= 0:
        return

    # Check if it's a picture
    if shape.shape_type == 13:  # PICTURE
        try:
            image = shape.image
            img_bytes = image.blob
            content_type = image.content_type
            pil_img = Image.open(io.BytesIO(img_bytes))
            pil_img = pil_img.convert('RGB')
            # Resize
            pil_img = pil_img.resize((w, h), Image.LANCZOS)
            img.paste(pil_img, (x, y))
        except Exception as e:
            pass
        return

    # Get fill color
    fill_color = None
    try:
        if hasattr(shape, 'fill') and shape.fill.type is not None:
            if shape.fill.type == 1:  # solid
                fill_color = get_color(shape.fill.fore_color.rgb)
    except:
        pass

    # Get line color
    line_color = None
    try:
        if hasattr(shape, 'line') and shape.line.color and shape.line.color.type == 1:
            line_color = get_color(shape.line.color.rgb)
    except:
        pass

    # Draw based on shape type
    st = shape.shape_type
    if st == 1:  # AUTO_SHAPE / rectangle
        # Check for rounded
        try:
            if hasattr(shape, 'adjustments') and shape.adjustments and shape.adjustments[0] > 0:
                # Rounded rectangle
                draw.rounded_rectangle([x, y, x + w, y + h], radius=int(min(w, h) * shape.adjustments[0]),
                                      fill=fill_color, outline=line_color, width=2)
            else:
                draw.rectangle([x, y, x + w, y + h], fill=fill_color, outline=line_color, width=2)
        except:
            draw.rectangle([x, y, x + w, y + h], fill=fill_color, outline=line_color, width=2)
    elif st == 9:  # OVAL
        draw.ellipse([x, y, x + w, y + h], fill=fill_color, outline=line_color, width=2)
    elif st == 17 or st == 19:  # LINE / CONNECTOR
        try:
            line_w = Pt(shape.line.width.pt if shape.line.width else 1)
            draw.line([(x, y), (x + w, y + h)], fill=fill_color or NAVY, width=int(line_w * 96 / 72))
        except:
            draw.line([(x, y), (x + w, y + h)], fill=fill_color or NAVY, width=2)
    else:
        # Default: rectangle
        if fill_color is not None or line_color is not None:
            try:
                if hasattr(shape, 'adjustments') and shape.adjustments and shape.adjustments[0] > 0:
                    draw.rounded_rectangle([x, y, x + w, y + h], radius=int(min(w, h) * shape.adjustments[0]),
                                          fill=fill_color, outline=line_color, width=2)
                else:
                    draw.rectangle([x, y, x + w, y + h], fill=fill_color, outline=line_color, width=2)
            except:
                draw.rectangle([x, y, x + w, y + h], fill=fill_color, outline=line_color, width=2)

    # Render text
    if shape.has_text_frame:
        render_text(draw, x, y, w, h, shape.text_frame)


def render_slide(slide, output_path, slide_w_emu, slide_h_emu, dpi=96):
    """Render a single slide to PNG."""
    w_px = emu_to_px(slide_w_emu, dpi)
    h_px = emu_to_px(slide_h_emu, dpi)
    img = Image.new('RGB', (w_px, h_px), WHITE)
    draw = ImageDraw.Draw(img)
    for shape in slide.shapes:
        try:
            render_shape(draw, shape, img)
        except Exception as e:
            pass
    img.save(output_path, optimize=True)


def main():
    pptx_path = '/home/user/AUREA_Gold_Standard_of_Data_V1.0.pptx'
    out_dir = '/home/user/aurea-pptx-previews'
    os.makedirs(out_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    print(f'Loaded {len(prs.slides)} slides at {prs.slide_width/914400:.2f}x{prs.slide_height/914400:.2f}in')

    for i, slide in enumerate(prs.slides):
        out = f'{out_dir}/slide_{i+1:02d}.png'
        render_slide(slide, out, prs.slide_width, prs.slide_height)
        print(f'  ✓ Rendered slide {i+1} → {out}')


if __name__ == '__main__':
    main()
